import io
from pptx import Presentation
from .base import BaseProcessor, ExtractionResult


class PptxProcessor(BaseProcessor):
    async def extract(self, file_content: bytes, filename: str) -> ExtractionResult:
        prs = Presentation(io.BytesIO(file_content))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text for cell in row.cells])
                    if rows:
                        slide_parts.append(
                            "Table: " + " | ".join(rows[0])
                        )
                        for row in rows[1:]:
                            slide_parts.append("  " + " | ".join(row))
            slides_text.append("\n".join(slide_parts))

        return ExtractionResult(
            text="\n\n".join(slides_text),
            tables=[],
            metadata={
                "filename": filename,
                "slide_count": len(prs.slides),
            },
        )
